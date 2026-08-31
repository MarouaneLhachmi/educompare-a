"""
Lecture des documents deposes — strategie par format.
======================================================

Jusqu'ici le systeme ne savait lire qu'un PDF pourvu d'une couche texte. Un
enseignant depose pourtant aussi bien un `.docx` qu'un support de
presentation, et un cours photocopie puis scanne n'a aucune couche texte.

Ce module centralise la lecture, pour deux raisons :

- l'Agent 1 et le controle au depot doivent lire **le meme document de la
  meme facon** ; les separer garantissait qu'ils divergent un jour ;
- les dependances de lecture sont optionnelles et leur absence doit produire
  un message exploitable, pas une trace technique. Il n'y a qu'un endroit ou
  le dire.

Strategie, dans l'ordre du plan de transition (phase 2.2) :

    .pdf    couche texte (pypdf)
            -> si vide : OCR (pdf2image + Tesseract), quand il est installe
    .docx   python-docx
    .pptx   python-pptx

Ce que les formats bureautiques apportent en plus
--------------------------------------------------
Un PDF ne dit pas ou commencent ses chapitres : il faut les deviner par
expressions regulieres. Un `.docx` porte ses styles de titre, un `.pptx` ses
titres de diapositive. Ces titres sont donc restitues separement, comme
**structure declaree par l'auteur**, et l'Agent 1 les prefere a sa propre
heuristique quand ils existent : une structure connue vaut mieux qu'une
structure devinee.

Pagination
----------
Un `.docx` n'a pas de pages : sa pagination depend de la police et de
l'imprimante. Plutot que d'inventer des numeros — qui se retrouveraient dans
l'annexe d'accreditation comme preuve localisee — le document est restitue en
une seule page et `pagination_fiable` vaut False. Un `.pptx`, lui, a des
diapositives : ce sont de vraies unites, numerotees comme telles.
"""

import os

SEUIL_TEXTE_EXPLOITABLE = 120  # caracteres en dessous desquels on tente l'OCR

FORMATS_SUPPORTES = ("pdf", "docx", "pptx")

LIBELLES_FORMAT = {
    "pdf": "PDF",
    "docx": "document Word (.docx)",
    "pptx": "présentation PowerPoint (.pptx)",
}


class DocumentIllisible(Exception):
    """Le document ne peut pas etre lu, avec un motif affichable."""


# ---------------------------------------------------------------------------
# Disponibilite des dependances optionnelles
# ---------------------------------------------------------------------------

def _importable(nom: str) -> bool:
    try:
        __import__(nom)
        return True
    except Exception:
        return False


def dependances() -> dict:
    """
    Etat des lecteurs, restitue dans la supervision technique. L'OCR compte
    trois pieces dont une hors de Python (le binaire Tesseract) : on ne peut
    donc pas se contenter de tester l'import.
    """
    ocr_python = _importable("pytesseract") and _importable("pdf2image")
    ocr_binaire = False
    if ocr_python:
        try:
            import pytesseract  # type: ignore

            pytesseract.get_tesseract_version()
            ocr_binaire = True
        except Exception:
            ocr_binaire = False
    return {
        "pdf": {"disponible": _importable("pypdf"), "module": "pypdf"},
        "docx": {"disponible": _importable("docx"), "module": "python-docx"},
        "pptx": {"disponible": _importable("pptx"), "module": "python-pptx"},
        "ocr": {
            "disponible": ocr_python and ocr_binaire,
            "module": "pytesseract + pdf2image + Tesseract",
            "python_installe": ocr_python,
            "binaire_installe": ocr_binaire,
        },
    }


def ocr_disponible() -> bool:
    return dependances()["ocr"]["disponible"]


def format_de(chemin: str) -> str:
    """Format deduit de l'extension, en minuscules et sans le point."""
    return os.path.splitext(chemin or "")[1].lstrip(".").lower()


def format_supporte(chemin_ou_extension: str) -> bool:
    extension = (chemin_ou_extension or "").lstrip(".").lower()
    if "." in (chemin_ou_extension or ""):
        extension = format_de(chemin_ou_extension)
    return extension in FORMATS_SUPPORTES


# ---------------------------------------------------------------------------
# Lecteurs par format
# ---------------------------------------------------------------------------

def _lire_pdf(chemin: str, max_pages: int | None) -> dict:
    from pypdf import PdfReader

    lecteur = PdfReader(chemin)
    pages_source = lecteur.pages if max_pages is None else lecteur.pages[:max_pages]
    pages = []
    for page in pages_source:
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            pages.append("")

    meta = lecteur.metadata or {}
    return {
        "pages": pages,
        "nb_pages_document": len(lecteur.pages),
        "titres": [],
        "metadonnees": {
            "titre_pdf": (meta.get("/Title") or "").strip() or None,
            "auteur_pdf": (meta.get("/Author") or "").strip() or None,
            "producteur": (meta.get("/Producer") or "").strip() or None,
        },
        "methode": "couche texte native du PDF (pypdf)",
        "pagination_fiable": True,
    }


def _ocr_pdf(chemin: str) -> list[str]:
    """
    Reconnaissance optique de caracteres pour les PDF scannes. Retourne une
    liste vide si l'OCR n'est pas installe : l'appelant decide quoi en dire.
    """
    try:
        import pytesseract  # type: ignore
        from pdf2image import convert_from_path  # type: ignore
    except Exception:
        return []
    try:
        images = convert_from_path(chemin, dpi=200)
        return [pytesseract.image_to_string(image, lang="fra+eng") for image in images]
    except Exception:
        return []


# Styles portant un titre, en francais comme en anglais : Word localise les
# noms de style selon la langue de l'installation qui a produit le fichier.
_STYLES_TITRE = ("heading", "titre", "title")


def _lire_docx(chemin: str, max_pages: int | None) -> dict:
    try:
        import docx  # type: ignore
    except Exception as exc:  # pragma: no cover - depend de l'environnement
        raise DocumentIllisible(
            "La lecture des documents Word nécessite la bibliothèque "
            "python-docx (pip install python-docx)."
        ) from exc

    document = docx.Document(chemin)
    lignes: list[str] = []
    titres: list[str] = []

    for paragraphe in document.paragraphs:
        texte = (paragraphe.text or "").strip()
        if not texte:
            continue
        style = (getattr(paragraphe.style, "name", "") or "").lower()
        if any(marque in style for marque in _STYLES_TITRE):
            titres.append(texte)
        lignes.append(texte)

    # Les tableaux portent souvent les exercices et les grilles de criteres :
    # les ignorer ferait disparaitre une part du contenu pedagogique.
    for tableau in document.tables:
        for rangee in tableau.rows:
            cellules = [c.text.strip() for c in rangee.cells if c.text.strip()]
            if cellules:
                lignes.append(" | ".join(cellules))

    proprietes = document.core_properties
    return {
        # Un .docx n'a pas de pages : en inventer produirait de fausses
        # references dans l'annexe d'accreditation.
        "pages": ["\n".join(lignes)],
        "nb_pages_document": 1,
        "titres": titres,
        "metadonnees": {
            "titre_pdf": (proprietes.title or "").strip() or None,
            "auteur_pdf": (proprietes.author or "").strip() or None,
            "producteur": "Microsoft Word (.docx)",
        },
        "methode": "structure du document Word (python-docx)",
        "pagination_fiable": False,
    }


def _lire_pptx(chemin: str, max_pages: int | None) -> dict:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover - depend de l'environnement
        raise DocumentIllisible(
            "La lecture des présentations nécessite la bibliothèque "
            "python-pptx (pip install python-pptx)."
        ) from exc

    presentation = Presentation(chemin)
    diapositives = list(presentation.slides)
    retenues = diapositives if max_pages is None else diapositives[:max_pages]

    pages: list[str] = []
    titres: list[str] = []
    for diapositive in retenues:
        morceaux: list[str] = []
        titre = None
        try:
            forme_titre = diapositive.shapes.title
            titre = (forme_titre.text or "").strip() if forme_titre is not None else None
        except Exception:
            titre = None
        if titre:
            titres.append(titre)
            morceaux.append(titre)

        for forme in diapositive.shapes:
            if not getattr(forme, "has_text_frame", False):
                continue
            texte = (forme.text_frame.text or "").strip()
            if texte and texte != titre:
                morceaux.append(texte)

        # Les notes du presentateur portent souvent le commentaire pedagogique
        # que la diapositive se contente de resumer.
        try:
            if diapositive.has_notes_slide:
                notes = (diapositive.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    morceaux.append(notes)
        except Exception:
            pass

        pages.append("\n".join(morceaux))

    proprietes = presentation.core_properties
    return {
        "pages": pages,
        "nb_pages_document": len(diapositives),
        "titres": titres,
        "metadonnees": {
            "titre_pdf": (proprietes.title or "").strip() or None,
            "auteur_pdf": (proprietes.author or "").strip() or None,
            "producteur": "Microsoft PowerPoint (.pptx)",
        },
        "methode": "diapositives et notes (python-pptx)",
        "pagination_fiable": True,
    }


_LECTEURS = {"pdf": _lire_pdf, "docx": _lire_docx, "pptx": _lire_pptx}


# ---------------------------------------------------------------------------
# Point d'entree
# ---------------------------------------------------------------------------

def lire(chemin: str, max_pages: int | None = None, tenter_ocr: bool = True) -> dict:
    """
    Lit un document, quel que soit son format.

    Retourne `pages`, `titres` (structure declaree par l'auteur, vide pour un
    PDF), `metadonnees`, `methode`, `format`, `ocr_utilise` et
    `pagination_fiable`.

    Leve `DocumentIllisible` avec un message affichable quand le format n'est
    pas pris en charge, que le lecteur manque, ou qu'aucun texte n'a pu etre
    obtenu.
    """
    if not os.path.exists(chemin):
        raise FileNotFoundError(f"Fichier introuvable : {chemin}")

    extension = format_de(chemin)
    lecteur = _LECTEURS.get(extension)
    if lecteur is None:
        formats = ", ".join(f".{f}" for f in FORMATS_SUPPORTES)
        raise DocumentIllisible(
            f"Format « .{extension or '?'} » non pris en charge. "
            f"Formats acceptés : {formats}."
        )

    try:
        resultat = lecteur(chemin, max_pages)
    except DocumentIllisible:
        raise
    except Exception as exc:
        raise DocumentIllisible(
            f"Le document n'a pas pu être lu ({LIBELLES_FORMAT.get(extension, extension)}) : "
            f"{str(exc)[:160]}"
        ) from exc

    resultat["format"] = extension
    resultat["ocr_utilise"] = False

    texte = "\n".join(resultat["pages"])
    if extension == "pdf" and len(texte.strip()) < SEUIL_TEXTE_EXPLOITABLE and tenter_ocr:
        pages_ocr = _ocr_pdf(chemin)
        if pages_ocr and len("".join(pages_ocr).strip()) >= SEUIL_TEXTE_EXPLOITABLE:
            resultat["pages"] = pages_ocr
            resultat["methode"] = "reconnaissance optique de caractères (Tesseract)"
            resultat["ocr_utilise"] = True
            texte = "\n".join(pages_ocr)

    if not texte.strip():
        if extension == "pdf":
            raise DocumentIllisible(
                "Aucun texte n'a pu être extrait de ce PDF. Le document est "
                "probablement une image scannée : installez Tesseract OCR "
                "(+ pdf2image) pour activer la reconnaissance optique de "
                "caractères."
            )
        raise DocumentIllisible(
            f"Aucun texte n'a pu être extrait de ce {LIBELLES_FORMAT.get(extension, extension)}. "
            "Le document est peut-être vide, ou son contenu est constitué "
            "d'images sans texte."
        )

    return resultat
