"""
Generateur du corpus de reference des tests de non-regression.
===============================================================

Les tests d'ancrage ont besoin de documents aux caracteristiques *connues* :
un cours complet doit obtenir une couverture elevee, un plan de cours une
couverture faible, un CV doit etre ecarte par le triage documentaire, et un
PDF sans couche texte doit provoquer un echec explicite de l'Agent 1.

Ces documents sont **generes** plutot que collectes, pour trois raisons :

- ils sont reproductibles a l'identique sur n'importe quelle machine, ce qui
  est la condition meme d'un intervalle d'ancrage stable ;
- ils ne posent aucune question de droit d'auteur ni de donnee personnelle ;
- leur contenu est ecrit pour cibler precisement une propriete a tester (le
  quasi-doublon differe du cours complet par quelques mots seulement, le plan
  de cours reprend les memes titres sans aucune substance).

Usage :
    python tests/corpus_reference/generer_corpus.py

Le catalogue des documents produits est ecrit dans `catalogue.json`, lu
ensuite par les tests.
"""

import json
import os

from fpdf import FPDF
from fpdf.enums import XPos, YPos

DOSSIER = os.path.dirname(os.path.abspath(__file__))
CATALOGUE_PATH = os.path.join(DOSSIER, "catalogue.json")


# ---------------------------------------------------------------------------
# Contenus
# ---------------------------------------------------------------------------

COURS_MATHS_COMPLET = [
    (
        "Chapitre 1 : Les grands nombres",
        "Dans ce chapitre, l'eleve apprend a lire, ecrire et comparer des nombres entiers "
        "jusqu'au million puis jusqu'au milliard. On travaille la decomposition en unites, "
        "dizaines, centaines, milliers et millions, ainsi que la valeur de position de chaque "
        "chiffre. L'eleve doit savoir ranger une liste de nombres dans l'ordre croissant et "
        "decroissant, encadrer un nombre entre deux dizaines ou deux centaines, et arrondir "
        "a la dizaine, a la centaine ou au millier le plus proche. De nombreux exercices "
        "d'application sont proposes, avec correction detaillee.",
    ),
    (
        "Chapitre 2 : Les fractions",
        "Les fractions simples sont introduites a partir du partage de figures : demis, tiers, "
        "quarts, cinquiemes. L'eleve apprend a reperer une fraction sur une droite graduee, a "
        "comparer deux fractions de meme denominateur, a reconnaitre des fractions equivalentes "
        "et a simplifier une fraction. L'addition et la soustraction de fractions de meme "
        "denominateur sont travaillees, puis la conversion d'une fraction decimale en nombre "
        "decimal et en pourcentage. Des exercices de calcul de pourcentage d'une quantite "
        "sont proposes en fin de chapitre.",
    ),
    (
        "Chapitre 3 : Les nombres decimaux",
        "Ecriture, lecture et comparaison des nombres decimaux. Valeur de position des chiffres "
        "apres la virgule : dixiemes, centiemes, milliemes. L'eleve apprend a additionner et "
        "soustraire des nombres decimaux en alignant les virgules, a multiplier un decimal par "
        "10, 100 et 1000, et a arrondir un decimal a l'unite ou au dixieme. Des problemes de "
        "monnaie et de mesure servent de support d'application.",
    ),
    (
        "Chapitre 4 : Multiplication et division",
        "Revision des tables de multiplication, technique de la multiplication posee a deux "
        "puis trois chiffres. Introduction de la division euclidienne avec quotient et reste, "
        "puis de la division decimale. L'eleve apprend a reconnaitre les multiples et les "
        "diviseurs d'un nombre, a appliquer les criteres de divisibilite par 2, 3, 5 et 9, et "
        "a resoudre des problemes concrets de partage et de groupement. Le calcul mental est "
        "travaille regulierement.",
    ),
    (
        "Chapitre 5 : Geometrie et figures planes",
        "Reconnaissance et proprietes du carre, du rectangle, du losange et des differents "
        "triangles. Notions de droites paralleles et perpendiculaires, tracees a la regle et a "
        "l'equerre. Calcul du perimetre et de l'aire du carre et du rectangle. Reperage de la "
        "symetrie axiale et construction du symetrique d'une figure. Introduction des solides "
        "usuels : cube, pave droit, cylindre, avec leurs faces, aretes et sommets.",
    ),
    (
        "Chapitre 6 : Mesures et grandeurs",
        "Unites de longueur (mm, cm, dm, m, km), de masse (g, kg, tonne), de contenance (mL, L) "
        "et de duree (seconde, minute, heure, jour). L'eleve apprend a convertir d'une unite a "
        "l'autre a l'aide d'un tableau de conversion, a calculer une duree entre deux instants, "
        "et a resoudre des problemes de la vie courante impliquant plusieurs unites. Le calcul "
        "d'aire en cm2 et m2 est mis en relation avec le chapitre de geometrie.",
    ),
    (
        "Chapitre 7 : Organisation et gestion de donnees",
        "Lecture et construction de tableaux a double entree, de diagrammes en batons et de "
        "graphiques. L'eleve apprend a extraire une information d'un tableau, a calculer une "
        "moyenne simple, et a reconnaitre une situation de proportionnalite. La regle de trois "
        "est introduite sur des exemples de recettes et d'echelles. Des exercices de "
        "resolution de problemes a plusieurs etapes cloturent le chapitre.",
    ),
]

# Le meme cours, ampute de cinq chapitres sur sept : la couverture doit chuter
# nettement par rapport au cours complet, sans tomber a zero.
COURS_MATHS_PARTIEL = COURS_MATHS_COMPLET[:2]

# Quasi-doublon : meme cours, quelques mots changes. Doit etre detecte par
# MinHash/LSH, et produire une analyse tres proche de l'originale.
COURS_MATHS_DOUBLON = [
    (titre, contenu.replace("L'eleve apprend", "Les eleves apprennent")
                   .replace("Dans ce chapitre", "Au fil de ce chapitre"))
    for titre, contenu in COURS_MATHS_COMPLET
]

# Plan de cours : les memes intitules, aucune substance enseignee. C'est le
# cas limite qui a motive la separation probabilite de couverture / suffisance
# dans l'Agent 6 : les notions sont *evoquees*, pas *enseignees*.
PLAN_DE_COURS = [
    ("Progression annuelle", "Chapitre 1 : les grands nombres. Chapitre 2 : les fractions. "
                             "Chapitre 3 : les nombres decimaux. Chapitre 4 : multiplication et division. "
                             "Chapitre 5 : geometrie et figures planes. Chapitre 6 : mesures et grandeurs. "
                             "Chapitre 7 : organisation et gestion de donnees."),
    ("Modalites", "Chaque chapitre est traite en quatre seances d'une heure. Une evaluation "
                  "sommative est organisee a la fin de chaque periode. Le manuel de reference "
                  "est fourni par l'etablissement. Les eleves disposent d'un cahier d'exercices."),
    ("Calendrier", "Periode 1 : chapitres 1 et 2. Periode 2 : chapitres 3 et 4. Periode 3 : "
                   "chapitres 5 et 6. Periode 4 : chapitre 7 et revisions generales."),
]

COURS_SCIENCES = [
    (
        "Chapitre 1 : La matiere et ses etats",
        "Les trois etats de la matiere : solide, liquide et gazeux. L'eleve observe les "
        "changements d'etat de l'eau : fusion, solidification, vaporisation et condensation. "
        "Le role de la temperature est mis en evidence par des experiences simples. La notion "
        "de melange homogene et heterogene est introduite, ainsi que les techniques de "
        "separation : decantation, filtration et evaporation.",
    ),
    (
        "Chapitre 2 : Le vivant et son environnement",
        "Classification des etres vivants, distinction entre vegetaux et animaux. Les besoins "
        "nutritifs des plantes et le role de la lumiere. Les chaines alimentaires et les "
        "reseaux trophiques dans un ecosysteme. L'eleve apprend a decrire les relations entre "
        "les etres vivants et leur milieu, et a identifier l'impact des activites humaines sur "
        "la biodiversite.",
    ),
    (
        "Chapitre 3 : Energie et circuits electriques",
        "Les differentes formes d'energie et leurs sources renouvelables ou non. Realisation "
        "d'un circuit electrique simple avec pile, interrupteur et lampe. Distinction entre "
        "conducteur et isolant. Les regles de securite face au courant electrique. L'eleve "
        "realise des experiences pour identifier les conditions de fonctionnement d'un circuit "
        "en serie et en derivation.",
    ),
    (
        "Chapitre 4 : Le corps humain",
        "Fonctionnement de l'appareil digestif, de l'appareil respiratoire et de la circulation "
        "sanguine. Le role d'une alimentation equilibree et de l'activite physique. Notions "
        "d'hygiene et de prevention. L'eleve apprend a decrire le trajet des aliments et celui "
        "de l'air dans l'organisme a l'aide de schemas legendes.",
    ),
]

# Cours en anglais : verifie que la comparaison francais/anglais fonctionne
# (c'est precisement ce que le repli LSA degrade — l'ancrage doit le montrer).
COURS_ANGLAIS = [
    (
        "Chapter 1: Place value and large numbers",
        "Pupils read, write, order and compare numbers up to ten million. They determine the "
        "value of each digit and round any whole number to a required degree of accuracy. "
        "Negative numbers are introduced in context and pupils calculate intervals across "
        "zero. Practice exercises cover ordering, rounding and estimation.",
    ),
    (
        "Chapter 2: Fractions, decimals and percentages",
        "Pupils compare and order fractions, including fractions greater than one. They add "
        "and subtract fractions with different denominators using the concept of equivalent "
        "fractions, multiply pairs of proper fractions, and divide proper fractions by whole "
        "numbers. Pupils recall and use equivalences between simple fractions, decimals and "
        "percentages, and solve problems involving the calculation of percentages.",
    ),
    (
        "Chapter 3: Geometry - properties of shapes",
        "Pupils draw two dimensional shapes using given dimensions and angles, recognise and "
        "build simple three dimensional shapes including making nets, and compare and classify "
        "geometric shapes based on their properties and sizes. They find unknown angles in any "
        "triangles, quadrilaterals and regular polygons, and illustrate the parts of a circle "
        "including radius, diameter and circumference.",
    ),
    (
        "Chapter 4: Measurement and statistics",
        "Pupils convert between miles and kilometres, calculate the area of parallelograms and "
        "triangles, and calculate the volume of cubes and cuboids. They interpret and construct "
        "pie charts and line graphs and use these to solve problems, and calculate and interpret "
        "the mean as an average of a set of data.",
    ),
]

CV = [
    ("Curriculum vitae", "Karim Bennani, 28 ans, ingenieur en informatique. Adresse : 14 rue des "
                         "Orangers, Casablanca. Telephone et courriel disponibles sur demande."),
    ("Experience professionnelle", "2021-2024 : developpeur backend chez une societe de services "
                                   "numeriques, en charge de la conception d'interfaces de "
                                   "programmation et de la maintenance de bases de donnees. "
                                   "2019-2021 : stage puis contrat de developpeur web."),
    ("Formation", "Diplome d'ingenieur en genie logiciel obtenu en 2019. Baccalaureat scientifique "
                  "mention bien en 2014."),
    ("Competences", "Python, Java, SQL, gestion de projet agile, anglais courant, arabe et francais "
                    "langues maternelles. Permis de conduire categorie B."),
]

FACTURE = [
    ("Facture n FA-2024-0187", "Emise le 12 mars 2024, echeance le 12 avril 2024. Vendeur : "
                               "Fournitures Bureautiques SARL, identifiant fiscal 40218765. "
                               "Client : etablissement scolaire Al Manar."),
    ("Detail des prestations", "Ramettes de papier A4, quantite 40, prix unitaire 42,00 dirhams, "
                               "total 1680,00. Cartouches d'encre, quantite 12, prix unitaire "
                               "185,00 dirhams, total 2220,00. Classeurs, quantite 60, prix "
                               "unitaire 15,50 dirhams, total 930,00."),
    ("Recapitulatif", "Total hors taxes 4830,00 dirhams. Taxe sur la valeur ajoutee 20 pour cent, "
                      "soit 966,00 dirhams. Net a payer 5796,00 dirhams. Reglement par virement "
                      "bancaire sous trente jours. Penalites de retard applicables au taux legal."),
]

CONTRAT = [
    ("Contrat de prestation de services", "Entre les soussignes, ci-apres denommes le prestataire "
                                          "et le client, il a ete convenu et arrete ce qui suit."),
    ("Article 1 - Objet", "Le present contrat a pour objet de definir les conditions dans "
                          "lesquelles le prestataire assure la maintenance du parc informatique "
                          "du client, ainsi que l'assistance aux utilisateurs."),
    ("Article 2 - Duree", "Le contrat est conclu pour une duree de douze mois a compter de sa "
                          "signature, renouvelable par tacite reconduction sauf denonciation "
                          "par lettre recommandee trois mois avant l'echeance."),
    ("Article 3 - Responsabilite et resiliation", "La responsabilite du prestataire ne saurait "
                                                  "etre engagee en cas de force majeure. Chaque "
                                                  "partie peut resilier le contrat en cas de "
                                                  "manquement grave de l'autre partie, apres mise "
                                                  "en demeure restee sans effet pendant trente jours."),
]


# ---------------------------------------------------------------------------
# Fabrication des PDF
# ---------------------------------------------------------------------------

def _ecrire_pdf(chemin: str, titre: str, sections: list[tuple[str, str]],
                sous_titre: str = "") -> None:
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 16)
    pdf.multi_cell(0, 10, titre, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    if sous_titre:
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(0, 6, sous_titre, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
    pdf.ln(6)

    for entete, contenu in sections:
        pdf.set_font("Helvetica", "B", 13)
        pdf.multi_cell(0, 8, entete, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_font("Helvetica", "", 11)
        pdf.multi_cell(0, 6, contenu, new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.ln(4)

    pdf.output(chemin)


def _ecrire_pdf_sans_texte(chemin: str) -> None:
    """
    PDF de deux pages ne contenant aucune couche texte : uniquement des
    rectangles. Simule un document scanne du point de vue de l'Agent 1, sans
    dependre d'une image externe ni de Pillow.
    """
    pdf = FPDF()
    for _ in range(2):
        pdf.add_page()
        pdf.set_fill_color(210, 210, 210)
        for i in range(12):
            pdf.rect(20, 25 + i * 18, 170 - (i % 3) * 25, 9, style="F")
    pdf.output(chemin)


def _ecrire_docx(chemin: str, titre: str, sections: list[tuple[str, str]]) -> None:
    """
    Document Word portant ses **styles de titre**. C'est ce qui distingue ce
    document du PDF equivalent : la structure y est declaree par l'auteur, pas
    devinee par expression reguliere. Les intitules sont volontairement
    depourvus du mot « Chapitre » — l'heuristique du PDF les manquerait, et
    c'est precisement ce que le test doit montrer.
    """
    import docx

    document = docx.Document()
    document.add_heading(titre, level=0)
    for entete, contenu in sections:
        document.add_heading(entete, level=1)
        document.add_paragraph(contenu)
    # Un tableau : les exercices y sont souvent, et les ignorer perdrait du
    # contenu pedagogique.
    tableau = document.add_table(rows=2, cols=2)
    tableau.cell(0, 0).text = "Exercice"
    tableau.cell(0, 1).text = "Competence visee"
    tableau.cell(1, 0).text = "Convertir 3,5 km en metres"
    tableau.cell(1, 1).text = "Conversion des unites de longueur"
    document.save(chemin)


def _ecrire_pptx(chemin: str, titre: str, sections: list[tuple[str, str]]) -> None:
    """
    Presentation dont chaque diapositive porte un titre et un corps, plus une
    note du presentateur — souvent le seul endroit ou le raisonnement
    pedagogique est ecrit en toutes lettres.
    """
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    couverture = presentation.slides.add_slide(presentation.slide_layouts[0])
    couverture.shapes.title.text = titre
    couverture.placeholders[1].text = "Support de presentation genere pour les tests."

    for entete, contenu in sections:
        diapositive = presentation.slides.add_slide(presentation.slide_layouts[1])
        diapositive.shapes.title.text = entete
        cadre = diapositive.placeholders[1].text_frame
        cadre.text = contenu[:220]
        for phrase in contenu[220:440].split(". "):
            if phrase.strip():
                paragraphe = cadre.add_paragraph()
                paragraphe.text = phrase.strip()
                paragraphe.font.size = Pt(14)
        diapositive.notes_slide.notes_text_frame.text = (
            "Note du presentateur : " + contenu[:180]
        )
    presentation.save(chemin)


# Intitules sans le mot « Chapitre » : l'heuristique du PDF ne les repererait
# pas, seuls les styles de titre les revelent.
SECTIONS_BUREAUTIQUES = [
    ("Les grands nombres et leur ecriture",
     COURS_MATHS_COMPLET[0][1]),
    ("Fractions, decimaux et pourcentages",
     COURS_MATHS_COMPLET[1][1] + " " + COURS_MATHS_COMPLET[2][1]),
    ("Figures planes, perimetres et aires",
     COURS_MATHS_COMPLET[4][1]),
    ("Unites de mesure et conversions",
     COURS_MATHS_COMPLET[5][1]),
]


# ---------------------------------------------------------------------------
# Catalogue
# ---------------------------------------------------------------------------
# `nature` decrit ce que le document EST ; `attendu` decrit ce que le systeme
# doit en dire. Les tests s'appuient sur ces deux champs, jamais sur le nom
# de fichier.

DOCUMENTS = [
    {
        "fichier": "cours_maths_complet.pdf",
        "titre": "Cours de Mathematiques - Derniere annee du primaire",
        "sous_titre": "Support de cours complet, sept chapitres.",
        "sections": COURS_MATHS_COMPLET,
        "nature": "cours_complet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "haute",
        },
    },
    {
        "fichier": "cours_maths_partiel.pdf",
        "titre": "Cours de Mathematiques - Premier semestre",
        "sous_titre": "Support couvrant les deux premiers chapitres seulement.",
        "sections": COURS_MATHS_PARTIEL,
        "nature": "cours_partiel",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "basse",
        },
    },
    {
        "fichier": "cours_maths_doublon.pdf",
        "titre": "Cours de Mathematiques - Derniere annee du primaire",
        "sous_titre": "Reformulation mineure du support complet.",
        "sections": COURS_MATHS_DOUBLON,
        "nature": "quasi_doublon",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "doublon_de": "cours_maths_complet.pdf",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "haute",
        },
    },
    {
        "fichier": "plan_de_cours.pdf",
        "titre": "Progression annuelle de mathematiques",
        "sous_titre": "Document d'organisation pedagogique, sans contenu enseigne.",
        "sections": PLAN_DE_COURS,
        "nature": "plan_de_cours",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "basse",
        },
    },
    {
        "fichier": "cours_sciences.pdf",
        "titre": "Cours de Sciences - Derniere annee du primaire",
        "sous_titre": "Support de cours complet, quatre chapitres.",
        "sections": COURS_SCIENCES,
        "nature": "cours_complet",
        "matiere": "Sciences",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "haute",
        },
    },
    {
        "fichier": "cours_maths_anglais.pdf",
        "titre": "Mathematics - Upper Key Stage 2",
        "sous_titre": "Complete course material, four chapters.",
        "sections": COURS_ANGLAIS,
        "nature": "cours_complet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "en",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "haute",
        },
    },
    {
        "fichier": "hors_sujet_cv.pdf",
        "titre": "Curriculum vitae",
        "sous_titre": "",
        "sections": CV,
        "nature": "hors_sujet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": False,
            "couverture_relative": "basse",
        },
    },
    {
        "fichier": "hors_sujet_facture.pdf",
        "titre": "Facture",
        "sous_titre": "",
        "sections": FACTURE,
        "nature": "hors_sujet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": False,
            "couverture_relative": "basse",
        },
    },
    {
        "fichier": "hors_sujet_contrat.pdf",
        "titre": "Contrat de prestation de services",
        "sous_titre": "",
        "sections": CONTRAT,
        "nature": "hors_sujet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": False,
            "couverture_relative": "basse",
        },
    },
    {
        "fichier": "cours_maths.docx",
        "titre": "Cours de Mathematiques - Derniere annee du primaire",
        "sous_titre": "",
        "sections": SECTIONS_BUREAUTIQUES,
        "generateur": "docx",
        "nature": "cours_complet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "haute",
        },
    },
    {
        "fichier": "cours_maths.pptx",
        "titre": "Mathematiques - Derniere annee du primaire",
        "sous_titre": "",
        "sections": SECTIONS_BUREAUTIQUES,
        "generateur": "pptx",
        "nature": "cours_complet",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": "fr",
        "attendu": {
            "extraction_reussie": True,
            "triage_pedagogique": True,
            "couverture_relative": "haute",
        },
    },
    {
        "fichier": "scan_sans_texte.pdf",
        "titre": "",
        "sous_titre": "",
        "sections": None,  # genere par _ecrire_pdf_sans_texte
        "nature": "scan_sans_couche_texte",
        "matiere": "Mathématiques",
        "niveau": "Dernière année du primaire",
        "langue": None,
        "attendu": {
            "extraction_reussie": False,
            "triage_pedagogique": None,
            "couverture_relative": None,
        },
    },
]


def generer() -> list[dict]:
    """Ecrit les PDF et le catalogue. Retourne la liste des documents."""
    catalogue = []
    for document in DOCUMENTS:
        chemin = os.path.join(DOSSIER, document["fichier"])
        generateur = document.get("generateur", "pdf")
        if document["sections"] is None:
            _ecrire_pdf_sans_texte(chemin)
        elif generateur == "docx":
            _ecrire_docx(chemin, document["titre"], document["sections"])
        elif generateur == "pptx":
            _ecrire_pptx(chemin, document["titre"], document["sections"])
        else:
            _ecrire_pdf(chemin, document["titre"], document["sections"],
                        document["sous_titre"])
        entree = {c: v for c, v in document.items() if c != "sections"}
        entree["taille_octets"] = os.path.getsize(chemin)
        catalogue.append(entree)

    with open(CATALOGUE_PATH, "w", encoding="utf-8") as fichier:
        json.dump(catalogue, fichier, ensure_ascii=False, indent=2)
    return catalogue


def main() -> None:
    catalogue = generer()
    print(f"{len(catalogue)} documents ecrits dans {DOSSIER}")
    for entree in catalogue:
        print(f"  {entree['fichier']:<28} {entree['nature']:<24} "
              f"{entree['taille_octets'] // 1024} Ko")
    print(f"Catalogue : {CATALOGUE_PATH}")


if __name__ == "__main__":
    main()
